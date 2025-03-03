import copy
import os
from PIL import Image
import pptx
import typing

from opencsp.common.lib.render.lib.PowerpointImage import PowerpointImage
from opencsp.common.lib.render.lib.PowerpointText import PowerpointText
from opencsp.common.lib.render_control.RenderControlPowerpointSlide import RenderControlPowerpointSlide
import opencsp.common.lib.tool.file_tools as ft
import opencsp.common.lib.tool.log_tools as lt


class PowerpointSlide:
    """
    Generate a powerpoint slide. Once completed, it should be added to a powerpoint presentation.
    """

    def __init__(
        self,
        slide_control: RenderControlPowerpointSlide,
        images: list[PowerpointImage] = None,
        texts: list[PowerpointText] = None,
    ):
        """Generate a powerpoint slide. Once completed, it should be added to a powerpoint presentation.

        Args:
            slide_control (RenderControlPowerpointSlide): Controls how this slide is rendered.
            images (list[PowerpointImage]): Images that this slide should start with. Defaults to None.
            texts (list[PowerpointText]): Texts that this slide should start with. Defaults to None.
        """
        self.slide_control = slide_control
        self.images: list[PowerpointImage] = images if images != None else []
        self.texts: list[PowerpointText] = texts if texts != None else []
        self.title_text_idx = None

        for pps_image in self.images:
            pps_image.parent_slide = self

        for idx, pps_text in enumerate(self.texts):
            pps_text.parent_slide = self
            found_title = False
            if not found_title:
                if pps_text.is_title:
                    self.title_text_idx = idx
                    found_title = True
            else:
                self.texts[idx] = pps_text._replace(is_title=False)

    @classmethod
    def template_title(cls, title: str, authors: str, slide_control: RenderControlPowerpointSlide) -> "PowerpointSlide":
        raise NotImplementedError

    @classmethod
    def template_planning(cls, slide_control: RenderControlPowerpointSlide = None) -> "PowerpointSlide":
        raise NotImplementedError

    @classmethod
    def template_overview(cls, slide_control: RenderControlPowerpointSlide = None) -> "PowerpointSlide":
        raise NotImplementedError

    @classmethod
    def template_content_simple(cls, slide_control: RenderControlPowerpointSlide = None) -> "PowerpointSlide":
        """Information-containing slide with a bulleted description on the left."""
        raise NotImplementedError

    @staticmethod
    def _get_cell_dims(
        table_dims: tuple[float, float, float, float],
        inter_cell_buffer_hor: float,
        inter_cell_buffer_vert: float,
        nrows: int,
        ncols: int,
        row_idx: int,
        col_idx: int,
    ):
        if row_idx >= nrows or col_idx >= ncols:
            lt.error_and_raise(
                ValueError,
                f"Error: in PowerpointSlide._get_cell_dims: row_idx ({row_idx}) or col_idx ({col_idx}) outside row/col bounds ({nrows}/{ncols})!",
            )
        x, y, w, h = table_dims

        inter_cell_buffer_hor_tot = inter_cell_buffer_hor * (ncols - 1)
        inter_cell_buffer_vert_tot = inter_cell_buffer_vert * (nrows - 1)
        cell_width = (w - inter_cell_buffer_hor_tot) / ncols
        cell_height = (h - inter_cell_buffer_vert_tot) / nrows

        cx = x + (cell_width * col_idx) + (inter_cell_buffer_hor * col_idx)
        cy = y + (cell_height * row_idx) + (inter_cell_buffer_vert * row_idx)
        cw = cell_width
        ch = cell_height

        return cx, cy, cw, ch

    @classmethod
    def template_content_grid(
        cls, nrows=2, ncols=2, slide_control: RenderControlPowerpointSlide = None
    ) -> "PowerpointSlide":
        """Information-containing slide, split into a ncols by nrows grid.

        Images or text boxes can be added to specific points in the grid by
        specifying the index, which proceeds from left to right, top to bottom.
        For example, in a 3x2 grid (3 columns, 2 rows), the cell indicies would
        be::

            0 | 1 | 2
            3 | 4 | 5

        Any space that an image isn't assigned to is assumed to be reserved
        to a text box.
        """
        if slide_control == None:
            slide_control = RenderControlPowerpointSlide(inter_cell_buffer=0.35)
        inter_cell_buffer = slide_control.inter_cell_buffer

        # get the content area
        left, top = slide_control.text_location
        right = left
        bottom = slide_control.title_location[1]
        caption_size = PowerpointText.compute_height(slide_control.image_caption_size)
        width = slide_control.slide_size[0] - (left + right)
        height = slide_control.slide_size[1] - (top + bottom)
        content_dims = left, top, width, height
        get_cell_dims = lambda row_idx, col_idx: cls._get_cell_dims(
            content_dims, inter_cell_buffer, inter_cell_buffer + caption_size, nrows, ncols, row_idx, col_idx
        )

        # generate texts
        texts: list[PowerpointText] = []
        for row in range(nrows):
            for col in range(ncols):
                texts.append(PowerpointText(None, cell_dims=get_cell_dims(row, col)))

        # generate images metadata
        images: list[PowerpointImage] = []
        for row in range(nrows):
            for col in range(ncols):
                images.append(PowerpointImage(None, cell_dims=get_cell_dims(row, col)))

        return cls(slide_control, images, texts)

    def _next_empty_cell_index(self, which_shapes="any"):
        if which_shapes == "any":
            shape_groups = [self.images, self.texts]
        elif which_shapes == "images":
            shape_groups = [self.images]
        elif which_shapes == "texts":
            shape_groups = [self.texts]
        else:
            lt.error_and_raise(
                ValueError,
                f"Error: in PowerpointSlide._next_empty_cell_index(), programmer error invalid value '{which_shapes}' for which_shapes argument",
            )

        ncells = max([len(shapes) for shapes in shape_groups])

        for idx in range(ncells):
            found = False

            for shapes in shape_groups:
                if len(shapes) > idx and shapes[idx].has_val():
                    found = True
                    break

            if not found:
                return idx

        return idx

    def add_image(self, image: PowerpointImage | typing.Any, fit_or_stretch: str = None, index: int = -1):
        """Add an image to this slide.

        If this slide has predefined spaces still available, then fits this image into the first of those spaces.
        Which space is filled can be overridden by specifying the index value.

        Args:
            - image (PowerpointImage|any): The image to be added. Can be any type that PowerpointImage accepts.
                                           At the time of this writing, this includes str, ndarray, Image, and RenderControlFigureRecord.
            - index (int): The cell index (position) in the template into which to set this image,
                           or -1 for the next available cell. Defaults to -1.
        """
        # set this as the image's parent
        if isinstance(image, PowerpointImage):
            image: PowerpointImage = copy.copy(image)
        elif image is None:
            lt.error_and_raise(
                ValueError, "Error: in PowerpointSlide.add_image(), can't add a 'None' type image to a slide."
            )
        else:
            image: PowerpointImage = PowerpointImage(image)
        if fit_or_stretch != None:
            image.stretch = fit_or_stretch.lower() == "stretch"
        image.parent_slide = self

        # find the next cell index into which to slot this image
        if index < 0 or index >= len(self.images):
            next_cell_idx = self._next_empty_cell_index("images")
            if next_cell_idx < len(self.images):
                index = next_cell_idx

        # add the image
        if index < 0:
            # append the image to the end of the images list
            self.images.append(image)
            index = len(self.images) - 1

            if not image.has_dims():
                slide_width, slide_height = self.slide_control.slide_size
                slide_dims = 0, 0, slide_width, slide_height
                image.fit_to_cell_dimensions(slide_dims)

        else:
            # slot the image into a given cell
            old_image = self.images[index]
            if old_image.has_val():
                old_image.clear_tmp_render()

            if not image.has_dims():
                self.images[index] = image
                image.fit_to_cell_dimensions(old_image.cell_dims)

    def add_text(self, text: PowerpointText, index: int = -1, replace_or_shift="replace"):
        if replace_or_shift not in ["replace", "shift"]:
            lt.error_and_raise(ValueError, f'Invalid argument replace_or_shift="{replace_or_shift}"')

        # set this as the text's parent
        text = copy.copy(text)
        text.parent_slide = self

        # find the next cell index into which to slot this text
        if index < 0 or index >= len(self.images):
            next_cell_idx = self._next_empty_cell_index("texts")
            if next_cell_idx < len(self.images):
                index = next_cell_idx

        # add the text
        if index < 0:
            self.texts.append(text)
            index = len(self.texts) - 1
        else:
            if index < len(self.texts):
                if replace_or_shift == "replace":
                    old_text = self.texts[index]
                    self.texts[index] = text
                    if text.dims == None:
                        text.dims = old_text.dims
                else:  # replace_or_shift == "shift"
                    for i in range(len(self.texts) - 1, index - 1, -1):
                        if i == len(self.texts) - 1:
                            self.texts.append(None)
                        self.texts[i + 1] = self.texts[i]
                    self.texts[index] = text
            else:
                self.texts.append(text)
                index = len(self.texts) - 1

        # set the title text index pointer
        if text.is_title:
            self.title_text_idx = index

        return text

    def get_title_text(self):
        if self.title_text_idx != None:
            return self.texts[self.title_text_idx]
        return None

    def set_title(self, title: str | PowerpointText):
        # get the title text instance
        title_text = self.get_title_text()
        if isinstance(title, str):
            slide_dims = 0, 0, *self.slide_control.slide_size
            if title_text != None:
                dims = title_text.dims
                text = PowerpointText(title, dims=dims, cell_dims=slide_dims, is_title=True, parent_slide=self)
            else:
                title_x, title_y = self.slide_control.title_location
                dims = title_x, title_y, slide_dims[2], 1
                text = PowerpointText(title, dims=dims, cell_dims=slide_dims, is_title=True, parent_slide=self)
                text.compute_and_assign_height(self.slide_control.title_size)
        else:
            text: PowerpointText = title
            text.is_title = True

        # add the text instance
        if self.title_text_idx != None:
            new_title_text = self.add_text(text, self.title_text_idx, replace_or_shift="replace")
        else:
            new_title_text = self.add_text(text, 0, replace_or_shift="shift")

        # set all other texts to not be titles
        for other_text in self.texts:
            if other_text != new_title_text:
                other_text.is_title = False

    def get_non_title_texts(self):
        if self.title_text_idx != None:
            non_title_texts = self.texts[: self.title_text_idx] + self.texts[self.title_text_idx + 1 :]
            return non_title_texts
        return self.texts

    def set_index(self, slide_index: int):
        if self.slide_control.slide_index == slide_index:
            return

        self.clear_tmp_saved_images_files()
        self.clear_tmp_saved_text_files()
        self.slide_control.slide_index = slide_index

    def clean(self, slide):
        """Clean the given pptx slide (remove all existing shapes)."""
        while len(slide.shapes) > 0:
            slide.shapes.element.remove(slide.shapes[0].element)
        while len(slide.placeholders) > 0:
            slide.placeholders.element.remove(slide.placeholders[0].element)

    def align_text(self, shape, alignment: pptx.enum.text.PP_ALIGN):
        text_frame = shape.text_frame
        for pidx in range(len(text_frame.paragraphs)):
            p = text_frame.paragraphs[pidx]
            p.alignment = alignment

    def format_text(self, shape, size: int):
        text_frame = shape.text_frame
        for pidx in range(len(text_frame.paragraphs)):
            p = text_frame.paragraphs[pidx]
            for ridx in range(len(p.runs)):
                r = p.runs[ridx]
                font = r.font
                font.size = pptx.util.Pt(size)

    def clear_tmp_saved_images_files(self):
        """Removes the temporary saved image files for this slide.
        To remove all saved images, use PowerpointImage.clear_tmp_save_all()."""
        for image in self.images:
            image.clear_tmp_save()

        # also remove all possible images for this slide's index
        if self.slide_control.slide_index >= 0:
            pattern = PowerpointImage._get_save_dir_name_ext_pattern(self.slide_control.slide_index, for_glob=True)
            dir, name, ext = ft.path_components(pattern)
            ft.delete_files_in_directory(dir, name + ext, error_on_dir_not_exists=False)

    def clear_tmp_saved_text_files(self):
        """Removes the temporary saved text files for this slide.
        To remove all saved texts, use PowerpointText.clear_tmp_save_all()."""
        for text in self.texts:
            text.clear_tmp_save()

    def save(self):
        """Saves images and texts out to temporary files, as necessary."""
        for image in self.images:
            image.save()
        for text in self.texts:
            text.save()

    def save_and_bake(self):
        """Saves the images and texts to temporary files, as with save().
        This also reduces image size and frees them from memory (note that this happens normally,
        it just happens in render instead)."""
        self.save()
        for image in self.images:
            if image.has_val():
                image.reduce_size(self.slide_control.reduced_image_size_scale)
        self.free_image_memory()

    def free_image_memory(self):
        # For images that have been saved, release the handle to their non-file format to save on memory.
        for image in self.images:
            image.replace_with_save()

    def render(self, presentation, layout, tmp_render_path: str):
        """Generates a slide, and adds all images and texts to be rendered to the slide.
        This has the side effect of calling save_and_bake()."""
        lt.info("In PowerpointSlide.render()")
        import opencsp.common.lib.render_control.RenderControlPowerpointPresentation as rcpp  # import here to avoid import loops

        presentation: rcpp.RenderControlPowerpointPresentation = presentation
        prs = presentation.presentation

        # get images ready to be added to the pp slide
        self.save_and_bake()

        # generate the slide
        prs.slide_width = pptx.util.Inches(self.slide_control.slide_size[0])
        prs.slide_height = pptx.util.Inches(self.slide_control.slide_size[1])
        slide = prs.slides.add_slide(layout)
        if presentation.existing_presentation_path_name_ext == None:
            self.clean(slide)

        # add all the images
        for image_idx, image in enumerate(self.images):
            # add the image to the slide
            if image.has_val():
                slide.shapes.add_picture(image.get_saved_path(), *image.dims_pptx())

            # add the caption, if any
            if image.caption != None:
                x, y, w, image_height_in = image.dims
                caption_size = 0.36  # TODO make this dependent on font point size
                h = caption_size
                if image.caption_is_above:
                    y -= caption_size
                else:
                    y += image_height_in
                caption_dims = [pptx.util.Inches(dim) for dim in [x, y, w, h]]
                txBox = slide.shapes.add_textbox(*caption_dims)
                txBox.text_frame.text = image.caption
                self.align_text(txBox, pptx.enum.text.PP_ALIGN.CENTER)
                self.format_text(txBox, self.slide_control.image_caption_size)

        # get the texts, adding the title text for the title slide
        if self.slide_control.is_title_slide:
            title_text = self.get_title_text()
            non_title_texts = self.get_non_title_texts()

            # add the title and subtitle to the slide
            if presentation.existing_presentation_path_name_ext != None:
                if title_text != None:
                    title = slide.shapes.title
                    title.text = title_text.val
                if len(non_title_texts) > 0:
                    subtitle = slide.placeholders[1]
                    subtitle.text = non_title_texts[0].val
                    non_title_texts = non_title_texts[1:]

            texts = non_title_texts
        else:
            texts = self.texts

        # add other texts to the slide
        for pps_text in texts:
            if pps_text.has_val():
                if not pps_text.has_dims():
                    lt.warn(
                        f'Warning: text "{pps_text.get_val()}" does not have a location set!\n\tOrigin: {pps_text.code_location}'
                    )
                    pps_text.dims = pps_text.cell_dims
                txBox = slide.shapes.add_textbox(*pps_text.dims_pptx())
                txBox.text_frame.text = pps_text.get_val()
                if pps_text.is_title:
                    self.format_text(txBox, self.slide_control.title_size)
                else:
                    self.format_text(txBox, self.slide_control.text_size)

        return slide

    def to_txt_file(self, file_path_name_ext: str):
        """Saves the images and texts for this slide in the given path
        and saves the references to them in the given file."""
        path, _, _ = ft.path_components(file_path_name_ext)

        # clean out existing files in the save directory
        if self.slide_control.slide_index >= 0:
            ft.delete_files_in_directory(path, f"{self.slide_control.slide_index}_*")

        # save images and texts to the save directory
        for image in self.images:
            image.set_save_path(path)
        for text in self.texts:
            text.set_save_path(path)
        self.save()

        # save references to the images and texts to the save directory
        with open(file_path_name_ext, "w") as fout:
            non_null = lambda v: v != None
            image_name_exts = list(filter(non_null, [image.save() for image in self.images]))
            text_name_exts = list(filter(non_null, [text.save() for text in self.texts]))

            fout.write("PowerpointSlide\n")
            fout.write("v1\n")
            fout.write(str(len(image_name_exts)) + "\n")
            fout.write(str(len(text_name_exts)) + "\n")
            for image_name_ext in image_name_exts:
                fout.write(image_name_ext + "\n")
            for text_name_ext in text_name_exts:
                fout.write(text_name_ext + "\n")

    @classmethod
    def from_txt_file(cls, file_path_name_ext: str, slide_control: RenderControlPowerpointSlide = None):
        if slide_control == None:
            slide_control = RenderControlPowerpointSlide()

        lines = ft.read_text_file(file_path_name_ext)
        lines = [line.strip() for line in lines]
        if len(lines) < 2:
            lt.error_and_raise(
                RuntimeError,
                f"Error: in PowerpointSlide.from_txt_file(), not enough lines in file {file_path_name_ext}",
            )

        file_type = lines[0]
        if file_type != "PowerpointSlide":
            lt.error_and_raise(
                RuntimeError,
                f"Error: in PowerpointSlide.from_txt_file(), bad file type {file_type} in {file_path_name_ext}, expected type PowerpointSlide",
            )
        version = lines[1]
        if version != "v1":
            lt.error_and_raise(
                RuntimeError,
                f"Error: in PowerpointSlide.from_txt_file(), bad version {version} in {file_path_name_ext}, expected version v1",
            )
        nimages = int(lines[2])
        ntexts = int(lines[3])

        images: list[PowerpointImage] = []
        texts: list[PowerpointText] = []
        for i in range(nimages):
            images.append(PowerpointImage.from_txt_file(lines[4 + i]))
        for j in range(ntexts):
            texts.append(PowerpointText.from_txt_file(lines[5 + i + j]))

        return cls(slide_control, images, texts)
