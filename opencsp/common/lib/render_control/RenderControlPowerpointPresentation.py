import os
import pptx

import opencsp.common.lib.opencsp_path.opencsp_root_path as orp
import opencsp.common.lib.render.PowerpointSlide as pps
import opencsp.common.lib.tool.file_tools as ft
import opencsp.common.lib.tool.log_tools as lt
import opencsp.common.lib.tool.time_date_tools as tdt


class RenderControlPowerpointPresentation:
    def __init__(self, new_slides: list[pps.PowerpointSlide] = None, existing_presentation_path_name_ext: str = None):
        """Create a new presentation instance to which to add slides.

        Example::

            # create the presentation to add our slides to
            presentation = rcpp.RenderControlPowerpointPresentation()

            # create each slide
            slide = pp.PowerpointSlide.template_content_grid(nrows=2, ncols=3)
            slide.add_image(top_left_image)
            slide.add_image(top_center_image)
            slide.add_image(top_right_image)
            slide.add_image(bottom_left_image)
            slide.add_image(bottom_center_image)
            slide.add_image(bottom_right_image)

            # add the slides
            slide.save_and_bake()
            presentation.add_slide(slide)

            # save the presentation
            presentation.save(home_dir() + "example.pptx")

        Args:
            new_slides (list[pps.PowerpointSlide]): Some of the slides to be added to the presentation. Defaults to None.
            existing_presentation_path_name_ext (str): An already existing presentation to which we want to add additional slides. Defaults to None.
        """
        self.new_slides = new_slides if new_slides != None else []
        self.existing_presentation_path_name_ext = existing_presentation_path_name_ext
        if existing_presentation_path_name_ext != None:
            lt.error_and_raise(
                NotImplementedError,
                "Error: in RenderControlPowerpointPresentation, support for use of existing presentations as a template not yet tested",
            )
        self.presentation = pptx.Presentation()

    def get_title_layout(self):
        return self.presentation.slide_layouts[0]

    def get_content_layout(self):
        return self.presentation.slide_layouts[1]

    def add_slide(self, slide: pps.PowerpointSlide):
        """Adds the given slide to this presentation and saves the associated images to temporary files.

        For this reason (saving images), adding the slide to the presentation
        should happen after the slide has been fully populated with contents."""
        self.new_slides.append(slide)
        slide.save()

    def save(self, dest_path_name_ext: str, overwrite=False):
        # check if the file already exists
        if ft.file_exists(dest_path_name_ext):
            if not overwrite:
                lt.error_and_raise(
                    FileExistsError,
                    f"Error: in RenderControlPowerpointPresentation.save: trying to save presentation to \"{dest_path_name_ext}\" but the file already exists",
                )

        dest_path, _, _ = ft.path_components(dest_path_name_ext)
        # check if the directory exists
        if not ft.directory_exists(dest_path):
            lt.error_and_raise(
                FileNotFoundError,
                f"Error: in RenderControlPowerpointPresentation.save: destination directory \"{dest_path}\" doesn't exist",
            )

        # setup
        tmp_dir = os.path.join(
            orp.opencsp_temporary_dir(), "powerpoint_presentations", tdt.current_date_time_string_forfile()
        )

        # render all the slides
        for slide_list_idx, pps_slide in enumerate(self.new_slides):
            control = pps_slide.slide_control
            layout = self.get_title_layout() if control.is_title_slide else self.get_content_layout()

            if control.slide_index < 0 or control.slide_index == slide_list_idx:
                pass
            else:
                raise NotImplementedError(
                    "Insertion of slides at a specific index hasn't been implemented in python-pptx (as of 06/28/2023)"
                )

            tmp_render_path = os.path.join(tmp_dir, f"slide_{slide_list_idx}")
            if ft.directory_exists(tmp_render_path):
                lt.error_and_raise(FileExistsError, f"Temporary rendering directory {tmp_render_path} already exists!")
            slide = pps_slide.render(self, layout, tmp_render_path)

        # check if the file already exists
        if ft.file_exists(dest_path_name_ext):
            # check again, just to be safe
            if not overwrite:
                lt.error_and_raise(
                    RuntimeError, f"Reached unreachable code!!! (file {dest_path_name_ext} already exists)"
                )

            # save to a temporary file
            name_ext = ft.body_ext_given_file_dir_body_ext(dest_path_name_ext)
            tmp_path_name_ext = os.path.join(tmp_dir, name_ext)
            ft.create_directories_if_necessary(tmp_dir)
            self.presentation.save(tmp_path_name_ext)

            # move to replace the existing file
            ft.delete_file(dest_path_name_ext)
            ft.rename_file(tmp_path_name_ext, dest_path_name_ext)
        else:
            self.presentation.save(dest_path_name_ext)

    @staticmethod
    def clear_tmp():
        pps.PowerpointImage.clear_tmp_save_all()
        pps.PowerpointText.clear_tmp_save_all()
